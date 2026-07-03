"""Build the Precedent Demo Day deck (submission + stage) as .pptx.

Verbatim transcription of Idea/refinement/03-pitch-deck.md, design system from
N1-deck-build.md §3, real numbers from Prep/final-numbers.md. Degraded rule applied to
team-names + Fetch URLs (no invented names, no ‹…›/[[WAIT]] tokens in the export).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0B, 0x0F, 0x14)
BODY = RGBColor(0xE6, 0xE6, 0xE6)
ACCENT = RGBColor(0x34, 0xD3, 0x99)
GREY = RGBColor(0x8A, 0x8F, 0x98)
DIM = RGBColor(0xB8, 0xBE, 0xC6)
FONT = "Arial"        # Inter is unavailable to soffice; Arial is the safe fallback
SW, SH = Inches(13.333), Inches(7.5)


def new_deck():
    p = Presentation()
    p.slide_width = SW
    p.slide_height = SH
    return p


def _blank(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def tb(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=None):
    """runs: list of paragraphs, each a list of (text, size, color, bold)."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space is not None:
            p.space_after = Pt(space)
        for (text, size, color, bold) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
    return box


def bar(s, x, y, w, h, color):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()
    r.shadow.inherit = False
    return r


def narr(s, text):
    """Self-narration band (grey, bottom) — submission deck only."""
    tb(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.8),
       [[(text, 12.5, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def build(with_band, with_monday):
    p = new_deck()

    # ---- Slide 1 — Title ----
    s = _blank(p)
    tb(s, Inches(1), Inches(2.4), Inches(11.3), Inches(2.2), [
        [("Precedent", 80, BODY, True)],
        [("Every incident resolved becomes precedent.", 26, ACCENT, False)],
    ], align=PP_ALIGN.CENTER)
    tb(s, Inches(1), Inches(6.2), Inches(11.3), Inches(0.6), [[
        ("Team — UK AI Agent Hackathon EP5  ·  Conduct “Make Legacy Move”  ·  Fetch.ai  ·  BasedAI",
         13, GREY, False)]], align=PP_ALIGN.CENTER)
    notes(s, "Say nothing on this slide — go straight into the cold open as slide 2 appears. "
             "Clicker in hand, demo machine mirrored, video cued.")

    # ---- Slide 2 — Cold open ----
    s = _blank(p)
    tb(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(1), [[
        ("Something breaks → find the manual → click through an admin console.", 24, BODY, True)]])
    tb(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.6), [
        [("8.85 hours", 88, ACCENT, True)],
        [("average time to resolve an incident — for a fix that is often minutes of "
          "keystrokes once it’s known", 20, DIM, False)],
    ])
    if with_band:
        narr(s, "When Sky News went off air in the CrowdStrike outage, the fix was a documented admin "
                "procedure executed thousands of times by hand — 8.85 business hours of queueing and "
                "lookup, not fixing (MetricNet/HDI).")
    notes(s, "Cold open, verbatim: 'When Sky News went off air in the CrowdStrike outage, the fix was a "
             "documented admin procedure — executed thousands of times, by hand. One of us watched that "
             "same loop every day inside Disney+ operations.' 8.85 business hours = MetricNet/HDI average "
             "MTTR. Do NOT claim Disney runs any particular vendor — first-hand observation only.")

    # ---- Slide 3 — Universal & priced ----
    s = _blank(p)
    tb(s, Inches(0.8), Inches(0.9), Inches(11.7), Inches(2.4), [
        [("$600B / year", 76, ACCENT, True)],
        [("downtime cost across the Global 2000 — up 50% in two years (Splunk, 2026)", 19, DIM, False)],
        [("≈ $200M per company (Oxford Economics/Splunk, 2024)", 22, BODY, True)],
    ])
    tb(s, Inches(0.8), Inches(4.9), Inches(11.7), Inches(1.4), [[
        (">60% of incidents are repeats — the fix already exists, nobody can find it "
         "(ServiceNow’s own support org)", 26, BODY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    if with_band:
        narr(s, "Downtime costs the Global 2000 ~$600B/yr — about $200M per company — and at "
                "ServiceNow’s own support desk, >60% of incidents were repeats whose fix already "
                "existed but couldn’t be found.")
    notes(s, "Pair the two numbers in ONE sentence. Sources verified 3–0: Splunk 2026; Oxford "
             "Economics/Splunk 2024 $200M/company; ServiceNow KCS case study (>60% repeats). Never say "
             "'$400B' bare — always pair $600B with $200M/company.")

    # ---- Slide 4 — Solution / thesis ----
    s = _blank(p)
    tb(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(1.6), [[
        ("AI SREs fix broken code. In real enterprises, the fix is almost never code — "
         "it’s a documented admin change.", 26, BODY, True)]])
    tb(s, Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.7), [
        [("Precedent remembers every fix your organisation has ever applied — and applies it again.",
          28, ACCENT, True)],
        [("risk-classified  ·  approval-gated  ·  audited  ·  reversible", 18, BODY, False)],
    ])
    tb(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.7), [[
        ("detect → find the documented fix → check risk → get approval → execute "
         "→ verify → remember", 15, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    if with_band:
        narr(s, "The memory unit is not a document — it’s an executed fix with provenance "
                "(symptom, verified fix, approver, risk class, rollback, outcome). A human approval gate, a "
                "full audit trail, and a rollback prepared before anything runs.")
    notes(s, "'Precedent is the agent that remembers every fix your organisation has ever applied and "
             "applies it again — with a human approval gate, a full audit trail, and a rollback prepared "
             "before anything runs.' Conduct framing: 'Conduct makes legacy legible; Precedent makes legacy "
             "operable.'")

    # ---- Slide 5 — Before/After stopwatch (time bars to scale) ----
    s = _blank(p)
    tb(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6), [[("Same incident. Same fix.", 24, BODY, True)]])
    full = Inches(11.7)
    tb(s, Inches(0.8), Inches(1.5), full, Inches(0.4),
       [[("Manual: ticket → search KB → admin console → approval queue → resolve",
          14, DIM, False)]])
    bar(s, Inches(0.8), Inches(1.95), full, Inches(0.55), RGBColor(0x3A, 0x41, 0x4B))
    tb(s, Inches(0.9), Inches(1.98), full, Inches(0.5), [[("8.85 hrs", 22, BODY, True)]],
       anchor=MSO_ANCHOR.MIDDLE)
    tb(s, Inches(0.8), Inches(3.0), full, Inches(0.35),
       [[("First time with Precedent (human approves)  ~60 s", 14, DIM, False)]])
    bar(s, Inches(0.8), Inches(3.4), Inches(0.5), Inches(0.5), ACCENT)
    tb(s, Inches(0.8), Inches(4.1), full, Inches(0.35),
       [[("Every time after (pre-approved)  ~15 s", 14, DIM, False)]])
    bar(s, Inches(0.8), Inches(4.5), Inches(0.16), Inches(0.5), ACCENT)
    tb(s, Inches(0.8), Inches(5.5), full, Inches(0.9), [[("The second time is free.", 40, ACCENT, True)]])
    if with_band:
        narr(s, "Three bars to scale: the manual loop is 8.85 hours of queueing and searching; Precedent’s "
                "first resolution ~60 s (a human still clicks approve); every repeat after ~15 s. The 15 s "
                "fast-path skips LLM triage for a known class — engineered, not hoped-for.")
    notes(s, "'The manual loop is eight and a half hours. Precedent’s first resolution ~a minute — a "
             "human clicks approve. Every repeat ~15 s. The second time is free.' This slide appears within the "
             "first 90 seconds. Never cut or degrade the time bars.")

    # ---- Slide 6 — The demo ----
    s = _blank(p)
    tb(s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.7), [[
        ("MediaCo — a simulated broadcaster  ·  real data  ·  live Jira Service Management  "
         "·  unscripted inputs", 18, BODY, True)]])
    tb(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(3.2), [
        [("1.  Broken EPG publish → plan + rollback → human approves → fixed in ~60 s, "
          "Jira ticket closes itself", 20, BODY, False)],
        [("2.  Same class recurs → pre-approved standard change → fixed in ~15 s", 20, ACCENT, True)],
        [("3.  Rights conflict, restricted runbook → Precedent refuses and routes to the right team",
          20, BODY, False)],
    ], space=16)
    tb(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.1), [[
        ("Seed data: real public runbooks (GitLab, Kubernetes SIG, the published CrowdStrike remediation), "
         "real programme metadata, and a 141k-event public IT-incident log ingested as fix history. Incident "
         "text is mutated at generation time — typos, vague symptoms, missing codes.", 11, GREY, False)]])
    if with_band:
        narr(s, "The incident text is deliberately mangled — typos, wrong terminology, missing codes — "
                "and any judge can file a ticket live. Incident 3: it isn’t allowed to read that runbook, so "
                "it refused and routed to the team that is. It knows what it’s not allowed to touch.")
    notes(s, "Say 'unscripted' on stage. Incident-3 payoff: 'It isn’t allowed to read that runbook — so "
             "it refused, and routed the incident to the team that is. It knows what it’s not allowed to "
             "touch.' The ~15 s is paced-up (real work ~6–8 s), never presented as raw latency.")

    # ---- Slide 7 — Control not autonomy ----
    s = _blank(p)
    tb(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8), [[
        ("Autonomy is earned per incident-class — and a human grants it.", 24, BODY, True)]])
    tb(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.9), [[
        ("L0 Observe  →  L1 Recommend  →  L2 Execute after approval  →  ", 20, DIM, False),
        ("L3 Standing Approval", 20, ACCENT, True),
        ("  (pre-approved standard change)", 16, DIM, False)]])
    b1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.3), Inches(3.3), Inches(0.75))
    b1.fill.solid(); b1.fill.fore_color.rgb = ACCENT; b1.line.fill.background(); b1.shadow.inherit = False
    b1.text_frame.text = "Promote to standing approval"
    for r in b1.text_frame.paragraphs[0].runs:
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = BG; r.font.name = FONT
    b2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.4), Inches(3.3), Inches(1.7), Inches(0.75))
    b2.fill.background(); b2.line.color.rgb = BODY; b2.line.width = Pt(1.25); b2.shadow.inherit = False
    b2.text_frame.text = "Revoke"
    for r in b2.text_frame.paragraphs[0].runs:
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = BODY; r.font.name = FONT
    tb(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.0), [[
        ("Approval moves earlier in time. It never leaves the loop.", 30, ACCENT, True)]])
    if with_band:
        narr(s, "No one approved that second ticket — the operations lead pre-approved the fix class after "
                "watching it succeed. Standing approval is an ITIL standard change: earned through verified "
                "history, revocable with one click, demoted automatically on any failure.")
    notes(s, "The permission decision is deterministic policy — no LLM ever authorises an action; the model "
             "proposes, the policy engine and a human identity dispose. ~8% of leaders accept full autonomy, "
             "~90% demand audit logs; the guardrails ARE the purchase criterion. L3 is never 'Autonomous'.")

    # ---- Slide 8 — How it works ----
    s = _blank(p)
    tb(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7), [[
        ("Five specialised agents  ·  deterministic policy engine  ·  permission-aware memory",
         20, BODY, True)]])
    # simple 5-node pipeline
    xs = [Inches(0.9 + i * 2.35) for i in range(5)]
    for i, x in enumerate(xs):
        n = s.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.4), Inches(1.7), Inches(1.1))
        n.fill.solid(); n.fill.fore_color.rgb = RGBColor(0x16, 0x1C, 0x24)
        n.line.color.rgb = ACCENT; n.line.width = Pt(1.25); n.shadow.inherit = False
        n.text_frame.text = ["detect", "triage", "decide", "execute", "verify"][i]
        for r in n.text_frame.paragraphs[0].runs:
            r.font.size = Pt(13); r.font.color.rgb = BODY; r.font.name = FONT
    shield = s.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(5.55), Inches(3.75), Inches(2.4), Inches(0.7))
    shield.fill.solid(); shield.fill.fore_color.rgb = ACCENT; shield.line.fill.background(); shield.shadow.inherit = False
    shield.text_frame.text = "deterministic policy gate — no LLM in the authorisation path"
    for r in shield.text_frame.paragraphs[0].runs:
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = BG; r.font.name = FONT
    cyl = s.shapes.add_shape(MSO_SHAPE.CAN, Inches(5.4), Inches(4.7), Inches(2.6), Inches(1.0))
    cyl.fill.solid(); cyl.fill.fore_color.rgb = RGBColor(0x16, 0x1C, 0x24); cyl.line.color.rgb = GREY; cyl.shadow.inherit = False
    cyl.text_frame.text = "executed fixes, permission-aware"
    for r in cyl.text_frame.paragraphs[0].runs:
        r.font.size = Pt(11); r.font.color.rgb = DIM; r.font.name = FONT
    tb(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.5), [[
        ("Live Jira Service Management  ·  open-weight models end-to-end  ·  agents live on "
         "Fetch.ai Agentverse / ASI:One", 14, GREY, False)]])
    if with_band:
        narr(s, "Five specialised agents run the loop; a deterministic policy engine — not the model — "
                "decides what may execute; and the memory inherits the permissions of the documents and tickets "
                "it learned from, live-synced from Jira. (Q&A-only on stage.)")
    notes(s, "100% open-weight: Qwen3.5-35B-A3B (Apache-2.0), DeepSeek-V4-Flash/Pro (MIT), BGE-M3 (MIT), Venice-"
             "served, verified vs public HF weights with /models dumps committed. Permission-aware memory ships "
             "as a standalone library Precedent imports. Three agents on Agentverse over Chat Protocol; the full "
             "loop runs inside an ASI:One conversation.")

    # ---- Slide 9 — Nobody does this ----
    s = _blank(p)
    tb(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.4), [[
        ("AI SREs stop at diagnosis. RPA is incident-blind. Runbook automation only runs pre-written "
         "scripts. ServiceNow deflects tickets inside its own walls.", 20, DIM, False)]])
    tb(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(2.0), [[
        ("The closed loop — retrieve your documented fix → approval-gated execution in third-party "
         "systems → verified, audited, remembered — is unclaimed.", 24, BODY, True)]])
    tb(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.0), [[
        ("ServiceNow paid $2.85B for Moveworks — resolution memory is worth buying.", 22, ACCENT, True)]])
    if with_band:
        narr(s, "$340M+ of AI-SRE VC since 2024 points at code/infra diagnosis and stops at 'here’s the "
                "likely cause.' Komodor/Ignio execute in their own silos — the whitespace is the "
                "combination: your documented fix + business-app execution + graduated approval + compounding "
                "memory. (Q&A-only on stage.)")
    notes(s, "Honest scoping: Komodor’s Klaudia and Digitate Ignio do execute remediations in their silos "
             "(K8s / their own automations) — the whitespace is the combination. Incumbents acquire this "
             "layer (Moveworks $2.85B, Mar 2025); that’s the exit narrative. Never say 'nobody executes'.")

    # ---- Slide 10 — The moat (REAL numbers baked in) ----
    s = _blank(p)
    tb(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.9), [[
        ("Every resolution writes a record: symptom → verified fix → approver → risk → "
         "rollback → outcome.", 20, BODY, True)]])
    tb(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.7), [[
        ("incident → fix → precedent → faster next time → more incidents trusted to it "
         "→ ↻", 18, ACCENT, True)]])
    tb(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(3.2), [
        [("141,000 real incident events ingested  ·  a 25k-record incident store", 18, BODY, True)],
        [("94% arrived with their fix already precedented — and still took a median 18 hrs "
          "(calendar) to resolve by hand", 18, BODY, False)],
        [("Permission-checked retrieval: P99 0.445 µs  ·  FNR 0 leaks / 5,219 deny-expected  "
          "·  6/6 adversarial attacks", 18, ACCENT, False)],
        [("Triage on mangled tickets (100-mutation corpus): 0 false-fast-paths / 100  ·  25/25 "
          "red-herring decoys resisted  (8% correct-match, 50% safe-degrade, 0% false fast-path)", 16, DIM, False)],
    ], space=12)
    if with_band:
        narr(s, "Documents go stale; executed fixes with provenance don’t. Every incident Precedent "
                "resolves makes the next one faster — accuracy becomes a function of tenure in the account. "
                "94% precedent existence is measured on 24,918 real UCI incidents; 18 hrs is calendar, never "
                "blended with the 8.85 business-hour baseline. (Q&A-only on stage.)")
    notes(s, "94% = fix-class EXISTENCE (key includes closed_code, known only at resolution); 98.6% at symptom "
             "level is the arrival-knowable number — don’t swap. P99 0.445 µs from bench/RESULTS.md "
             "(seed 4207). Extractor: 0 false-fast-paths / 100 mutations, 25/25 decoys resisted "
             "(extractor_robustness.json). 18 hrs is CALENDAR — keep the label.")

    # ---- Slide 11 — Market ----
    s = _blank(p)
    tb(s, Inches(0.8), Inches(0.8), Inches(11.7), Inches(3.4), [
        [("Media ops first — regulated failure (Ofcom), 24/7 deadlines, legacy consoles", 20, BODY, True)],
        [("One MSP deal = hundreds of channels — Encompass runs 1,200+ channels/day; Amagi 5,000+ "
          "deliveries (vendor-claimed)", 18, BODY, False)],
        [("Then: every ops team drowning in runbooks — every ticket deflected saves "
          "$22 → $69 → $104 per escalation tier", 18, ACCENT, True)],
    ], space=18)
    tb(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.0), [[
        ("Agents made retrieve-and-safely-execute possible this year. Incumbents are priced to meter the "
         "workflow, not delete it.", 16, GREY, False)]])
    if with_band:
        narr(s, "Media first: first-hand pain, regulator-documented failures, Netflix proved the pattern "
                "in-house. The MSP channel is the wedge-scaler — fixed-price SLA ops means every "
                "auto-remediated incident is pure margin. $22/$69/$104 = MetricNet cost-per-ticket "
                "(vintage ~2019–20).")
    notes(s, "Vendor-claimed marketing figures (Encompass/Amagi) — say so if pressed. $22/$69/$104 = "
             "MetricNet whitepaper, vintage ~2019–2020 — quote the numbers, not fake freshness.")

    # ---- Slide 12 — Team + Ask (degraded: no invented names) ----
    s = _blank(p)
    tb(s, Inches(0.8), Inches(0.9), Inches(11.7), Inches(1.2), [[
        ("Founding team — UK AI Agent Hackathon EP5", 22, BODY, True)]])
    tb(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.0), [[
        ("(one of us watched the incident loop first-hand inside broadcast-streaming operations)", 15, GREY, False)]])
    tb(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.7), [[
        ("The ask: two intros to broadcast-ops or MSP design partners — and we’re applying to "
         "Antler & EWOR with this.", 26, ACCENT, True)]])
    tb(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.0), [[
        ("ServiceNow paid $2.85B for resolution memory. We’re building the version that executes.",
         20, BODY, True)]])
    if with_band:
        narr(s, "The ask, as a sentence: two introductions to anyone running a 24/7 ops or NOC team — "
                "media is where we start, not where we stop — and we’re applying to Antler and EWOR. "
                "(Team names/photos are added by the founders before the live pitch.)")
    notes(s, "Deliver the ask as a sentence, not a slide-read. Team names/photos/credentials are human-written "
             "and added before the live pitch — never AI-invented. Last words before Q&A: 'The second time "
             "is free.'")

    # ================= APPENDIX (A1..A7, A9, A8) =================
    def appendix_title(s, t):
        tb(s, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.6), [[(t, 22, ACCENT, True)]])

    # A1 — Every number, sourced
    s = _blank(p); appendix_title(s, "A1 — Every number, sourced")
    a1 = [
        ("Global 2000 downtime bill", "$600B/yr, +50% in 2 yrs", "Splunk 2026 (exact framing)"),
        ("Per-company downtime cost", "$200M/yr", "Oxford Economics/Splunk 2024 — verified 3–0"),
        ("Repeat incidents with existing fixes", ">60%", "ServiceNow KCS case study — verified 3–0"),
        ("Cost per ticket L1 → desktop → L3", "$22 → $69 → $104", "MetricNet whitepaper (~2019–20)"),
        ("Average incident MTTR", "8.85 business hrs", "MetricNet/HDI"),
        ("KB-attached / KCS relief", "66% faster / 52% faster", "ServiceNow — verified"),
        ("Moveworks acquisition", "$2.85B", "ServiceNow newsroom, Mar 2025"),
        ("Deflection precedent", "50–88%", "Moveworks — vendor-claimed"),
        ("Fix-class match rate (24,918 real incidents)", "94.4% (symptom 98.6%)", "our measurement, UCI CC BY 4.0"),
        ("Median resolution, precedented repeats", "18.2 calendar hrs (p75 136.6)", "ours — never blend w/ 8.85 business"),
        ("Recurring classes ≥4 occurrences", "558 classes = 94.8% of volume", "ours — ladder-bootstrap evidence"),
        ("Extractor robustness (100-mutation corpus)", "0 false-fast-paths / 100; 25/25 decoys resisted", "ours — seed 4207"),
        ("Conformance FNR / attacks", "0 leaks / 5,219; 6/6", "ours — independent-oracle bench"),
    ]
    rows = []
    for (c, n, src) in a1:
        rows.append([(c + "  ", 12, BODY, False), (n + "  ", 12, ACCENT, True), (src, 11, GREY, False)])
    tb(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.9), rows, space=5)

    # A2 — Data provenance
    s = _blank(p); appendix_title(s, "A2 — Data provenance (the “not made up” slide)")
    tb(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.4), [
        [("KB seeded with real published runbooks (GitLab, Kubernetes SIG, the CrowdStrike channel-file "
          "remediation) — every article carries its adapted_from URL.", 15, BODY, False)],
        [("Fix history from the UCI ServiceNow incident event log (141,712 events / 24,918 incidents, "
          "CC BY 4.0) — measured 94% fix-class match, 18.2h calendar median for precedented repeats.", 15, BODY, False)],
        [("Scheduler/rights/EPG seeded with real UK programme metadata (TVmaze CC BY-SA / Freeview XMLTV) and "
          "real streaming-catalog data (CC0 Kaggle) — only licence-window terms are synthesised, by a "
          "stated rule.", 15, BODY, False)],
        [("We checked licences and REJECTED TMDB (API terms prohibit AI/ML use) and IMDb — the diligence "
          "is itself a credibility beat.", 15, ACCENT, True)],
        [("Incident text mutated at generation time (typos, colloquial symptoms, missing codes, red herrings). "
          "The systems are simulated; the content is real.", 15, DIM, False)],
    ], space=12)

    # A3 — Standing approval semantics
    s = _blank(p); appendix_title(s, "A3 — Standing approval: exact semantics")
    tb(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.4), [
        [("Incident-class key = deterministic fingerprint sha256(service | error_code | target_object_type) "
          "from structured fields — the LLM proposes candidates, but a match only counts on fingerprint "
          "equality. The model never authorises itself.", 15, BODY, False)],
        [("Graduation, printed in the audit log: 3 consecutive verified L2 successes, zero rollbacks → "
          "eligible; a human clicks Promote to Standing Approval. Any verification failure or rollback "
          "auto-demotes the class to L1 and logs a demotion event.", 15, BODY, False)],
        [("The rollback plan is generated BEFORE execution and is a precondition of the gate; verification "
          "failure fires it automatically.", 15, BODY, False)],
        [("Approver identity is recorded per action (chat: the Chat Protocol sender address; production: SSO "
          "identity — the control structure is identical).", 15, BODY, False)],
    ], space=14)

    # A4 — Permission-aware memory (BasedAI)
    s = _blank(p); appendix_title(s, "A4 — Permission-aware memory (BasedAI depth)")
    tb(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.4), [
        [("Every record stores the FULL set of source permission constraints; retrieval must satisfy ALL of "
          "them — conjunction, not one strictest label.", 15, BODY, False)],
        [("A precompiled effective-policy cache makes the check one indexed lookup: P99 0.445 µs (end-to-end "
          "overhead P99 0.0130 ms) over the 25k-record store, concurrent queries.", 15, ACCENT, True)],
        [("Benchmarked in the track’s own vocabulary: FNR 0 / 5,219 and FPR 0 / 4,781 over 10,000 "
          "ground-truth queries graded by an INDEPENDENT oracle, P50/P99, overhead, ACL drift, "
          "time-to-consistency — each vs threshold in bench/RESULTS.md; six named attacks pass.", 15, BODY, False)],
        [("Revocation cascades — dual enforcement: flip the issue security level in Jira → Jira hides "
          "the runbook AND every derived fix/summary/embedding becomes unretrievable within one poll tick, with "
          "denial audit events in BOTH logs.", 15, BODY, False)],
        [("Fallback mode fails closed: uncertain ACL freshness → not served. Ships as a standalone library; "
          "100% open-weight models, named in the README.", 15, DIM, False)],
    ], space=11)

    # A5 — Why won't ServiceNow build this
    s = _blank(p); appendix_title(s, "A5 — Why won’t ServiceNow / Conduct build this?")
    tb(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.4), [
        [("ServiceNow monetises seats + tickets; auto-resolution cannibalises both. They ACQUIRE this layer: "
          "Jeli $29.7M, Moveworks $2.85B. We’re building what they buy.", 15, BODY, False)],
        [("Now Assist acts inside ServiceNow workflows; we execute inside third-party business-app admin "
          "surfaces — their walls are our territory.", 15, BODY, False)],
        [("Conduct makes legacy legible; we make it operable — a complement, and the brief told us to build "
          "exactly this (“start with the process”).", 15, BODY, False)],
        [("Rehearsed: “They meter the workflow; we delete it. Incumbents don’t ship products that shrink "
          "their own ticket count — they buy them. That’s the Moveworks story.”", 15, ACCENT, True)],
    ], space=14)

    # A6 — Integration reality
    s = _blank(p); appendix_title(s, "A6 — Integration reality (the “you built the simulator” defence)")
    tb(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.4), [
        [("Four-tier surface strategy: REST APIs (WhatsOn-class schedulers) → BXF/file exchange → "
          "watched-folder/FTP drops → RPA-style console driving (last resort, same approval gates).", 15, BODY, False)],
        [("The Operator executes only TYPED tool calls per surface — never free-form shell; tier 3/4 fixes "
          "start at L0/L1 and stay approval-gated.", 15, BODY, False)],
        [("MediaCo’s endpoints deliberately mirror the surfaces the real vendors document (Grass Valley "
          "STRATUS: “traffic integration via BXF files”).", 15, BODY, False)],
    ], space=16)

    # A7 — Fetch.ai deliverables (degraded: no fabricated URLs)
    s = _blank(p); appendix_title(s, "A7 — Fetch.ai deliverables")
    tb(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.4), [
        [("Three agents registered on Agentverse as mailbox agents, collaborating over the uAgents Agent Chat "
          "Protocol; both Innovation Lab + hackathon badges on each profile; the hosted degraded-L0 Watcher "
          "stays live post-hackathon.", 15, BODY, False)],
        [("The full loop runs inside an ASI:One conversation: report an incident in chat → plan + rollback "
          "→ “approve” → the ticket transitions and closes → audit hash returned.", 15, BODY, False)],
        [("Approval principal = the Chat Protocol sender address, logged as the authorising identity, with a "
          "10-minute gate TTL (a dropped session never leaks an execution).", 15, BODY, False)],
        [("Agentverse profile URLs + the ASI:One shared-chat URL are captured at live registration.", 14, GREY, False)],
    ], space=14)

    # A9 (before A8) — Bottoms-up ACV
    s = _blank(p); appendix_title(s, "A9 — Bottoms-up ACV")
    tb(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.2), [
        [("One arithmetic slide, every input sourced:", 16, BODY, True)],
        [("1 MSP NOC (Encompass-class, 1,200+ channels)  ×  tens of thousands of tickets/yr  ×  60%+ "
          "repeat share (ServiceNow, verified 3–0; our corpus says 94% precedented)  ×  $50 blended "
          "deflection value (inside the $22–$104 MetricNet ladder)", 16, BODY, False)],
        [("⇒  seven-figure ACV justification per site — before counting downtime avoidance.", 20, ACCENT, True)],
        [("This is the bottoms-up serviceable-market answer; it does not replace primary customer evidence and "
          "never pretends to.", 14, GREY, False)],
    ], space=16)

    # A8 — Liability & regulated ops
    s = _blank(p); appendix_title(s, "A8 — Liability & regulated ops")
    tb(s, Inches(0.7), Inches(1.2), Inches(11.9), Inches(5.4), [
        [("Precedent executes only documented, previously-verified fixes at the customer’s chosen autonomy "
          "level — the customer’s own change-management policy, encoded and enforced. Shadow-mode-first "
          "at MSPs.", 15, BODY, False)],
        [("Headline metric we quote: per-class execution success + rollback rate from real runs — never "
          "model benchmarks (the incident.io-vs-Rootly accuracy-theatre war is the cautionary tale).", 15, BODY, False)],
        [("Every action carries: input, reasoning, confidence, action, authorising rule + approver — the "
          "audit schema regulators and SOC 2 auditors already recognise.", 15, BODY, False)],
    ], space=16)

    # ---- PDF-only: What exists Monday morning ----
    if with_monday:
        s = _blank(p); appendix_title(s, "What exists Monday morning")
        tb(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(5.0), [
            [("Durable artifacts — not slideware:", 18, BODY, True)],
            [("•  A hosted degraded-L0 Watcher, live on Fetch.ai Agentverse (answers a described incident "
              "at L0, never executes).", 16, BODY, False)],
            [("•  precedent_memory — the permission-aware memory as a standalone, importable Python "
              "library (fail-closed retrieval, hash-chained audit, conformance bench).", 16, BODY, False)],
            [("•  The ground-truth conformance bench, graded by an independent oracle (FNR 0 / 5,219, "
              "6/6 attacks, seed 4207 — byte-identically reproducible).", 16, BODY, False)],
            [("•  The public evidence index (docs/evidence, docs/compliance) — every claim traceable "
              "to a committed artifact.", 16, BODY, False)],
        ], space=14)

    return p


for name, band, monday in [("precedent-deck", True, True), ("precedent-deck-stage", False, False)]:
    build(band, monday).save(f"docs/deck/{name}.pptx")
    print(f"wrote docs/deck/{name}.pptx")
